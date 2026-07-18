#!/usr/bin/env python3
"""
Generador del deck "Claude Code — Curso en dos partes".

Una sola presentación con dos partes diferenciadas:
  Parte 1 — Claude Code (la herramienta)
  Parte 2 — La metodología (agnóstica de la herramienta)

Estilo: slides en blanco dibujadas a mano con eyebrow + título + tarjetas
numeradas + agenda + franja de conclusión + footer. Idempotente: re-ejecutar
regenera el .pptx.

    pip install python-pptx
    python build_pptx.py

Salida: Claude_Code_Presentacion.pptx (16:9, 13.33 x 7.5").
"""
from __future__ import annotations
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ----------------------------------------------------------------------------- paleta
BG       = RGBColor(0x0F, 0x14, 0x1A)   # slate casi negro
PANEL    = RGBColor(0x1A, 0x21, 0x2B)   # tarjeta
PANEL_2  = RGBColor(0x22, 0x2B, 0x38)   # tarjeta alt / highlight
STROKE   = RGBColor(0x2E, 0x39, 0x48)   # borde sutil
CORAL    = RGBColor(0xD9, 0x77, 0x57)   # acento primario (Anthropic)
BLUE     = RGBColor(0x7A, 0xA2, 0xC7)   # acento secundario
GREEN    = RGBColor(0x8C, 0xC2, 0x8C)
TEXT     = RGBColor(0xE8, 0xEA, 0xED)   # texto principal
MUTED    = RGBColor(0x9A, 0xA4, 0xB2)   # texto secundario
FAINT    = RGBColor(0x66, 0x70, 0x7E)   # footer
FONT     = "Segoe UI"
MONO     = "Consolas"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)


# ----------------------------------------------------------------------------- helpers
def _solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _noline(shape):
    shape.line.fill.background()


def _line(shape, color, w=0.75):
    shape.line.color.rgb = color
    shape.line.width = Pt(w)


def rect(slide, x, y, w, h, fill=None, line=None, lw=0.75, rounded=False, radius=0.08):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    if rounded:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        _solid(shp, fill)
    if line is None:
        _noline(shp)
    else:
        _line(shp, line, lw)
    shp.shadow.inherit = False
    return shp


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.0):
    """runs: lista de párrafos; cada párrafo es lista de (txt, size, color, bold, font)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, size, color, bold, font) in para:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = font
    return tb


def R(t, size, color, bold=False, font=FONT):
    return (t, size, color, bold, font)


# ----------------------------------------------------------------------------- chrome
def base(slide, eyebrow=None, ttl=None, page=None):
    rect(slide, -0.1, -0.1, 13.53, 7.7, fill=BG)                     # fondo
    rect(slide, 0, 0, 0.22, 7.5, fill=CORAL)                         # barra lateral
    if eyebrow:
        text(slide, 0.7, 0.55, 11, 0.4,
             [[R(eyebrow.upper(), 12, CORAL, True)]])
    if ttl:
        text(slide, 0.68, 0.92, 12, 1.1, [[R(ttl, 30, TEXT, True)]])
        rect(slide, 0.72, 1.72, 1.1, 0.045, fill=CORAL)
    if page is not None:
        text(slide, 0.7, 7.02, 8, 0.35,
             [[R("Claude Code · Curso / Workshop", 9.5, FAINT, False)]])
        text(slide, 11.8, 7.02, 1.2, 0.35,
             [[R(str(page), 9.5, FAINT, True)]], align=PP_ALIGN.RIGHT)


def takeaway(slide, phrase):
    """Franja de conclusión (🗣️) al pie."""
    y = 6.35
    rect(slide, 0.7, y, 11.95, 0.52, fill=PANEL_2, line=STROKE, rounded=True, radius=0.28)
    rect(slide, 0.7, y, 0.09, 0.52, fill=CORAL, rounded=True, radius=0.5)
    text(slide, 0.95, y + 0.03, 11.6, 0.46,
         [[R("“ ", 13, CORAL, True), R(phrase, 12.5, TEXT, False),
           R(" ”", 13, CORAL, True)]], anchor=MSO_ANCHOR.MIDDLE)


def cards(slide, items, top=2.1, cols=2, height=1.02, gap=0.28, badge=CORAL, bottom_takeaway=True):
    """items: lista de (num, titulo, descripcion). Rejilla de tarjetas numeradas."""
    left, right = 0.7, 12.65
    total_w = right - left
    card_w = (total_w - gap * (cols - 1)) / cols
    for idx, (num, ttl, desc) in enumerate(items):
        r, c = divmod(idx, cols)
        x = left + c * (card_w + gap)
        y = top + r * (height + gap)
        rect(slide, x, y, card_w, height, fill=PANEL, line=STROKE, rounded=True, radius=0.09)
        # badge
        rect(slide, x + 0.22, y + 0.22, 0.5, 0.5, fill=badge, rounded=True, radius=0.32)
        text(slide, x + 0.22, y + 0.24, 0.5, 0.46,
             [[R(str(num), 15, BG, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(slide, x + 0.9, y + 0.16, card_w - 1.05, 0.4,
             [[R(ttl, 14.5, TEXT, True)]])
        text(slide, x + 0.9, y + 0.5, card_w - 1.05, height - 0.55,
             [[R(desc, 11, MUTED, False)]], line_spacing=1.02)


def code(slide, x, y, w, h, lines, title=None):
    if title:
        text(slide, x, y - 0.34, w, 0.3, [[R(title, 11, BLUE, True)]])
    rect(slide, x, y, w, h, fill=RGBColor(0x0A, 0x0E, 0x13), line=STROKE, rounded=True, radius=0.05)
    runs = [[R(ln, 11.5, TEXT if not ln.startswith("#") else MUTED, False, MONO)] for ln in lines]
    text(slide, x + 0.28, y + 0.2, w - 0.5, h - 0.35, runs, line_spacing=1.12, space_after=2)


def panel_bullets(slide, x, y, w, h, header, bullets, accent=CORAL):
    rect(slide, x, y, w, h, fill=PANEL, line=STROKE, rounded=True, radius=0.06)
    text(slide, x + 0.3, y + 0.22, w - 0.5, 0.4, [[R(header, 14, accent, True)]])
    runs = []
    for b in bullets:
        runs.append([R("▸  ", 11.5, accent, True), R(b, 11.5, TEXT, False)])
    text(slide, x + 0.3, y + 0.72, w - 0.55, h - 0.9, runs, line_spacing=1.05, space_after=6)


def divider(slide, part, ttl, subtitle, items, page=None):
    """Separador de parte: número grande + título + lista de secciones."""
    rect(slide, -0.1, -0.1, 13.53, 7.7, fill=BG)
    rect(slide, 0, 0, 13.333, 0.16, fill=CORAL)
    rect(slide, 0, 7.34, 13.333, 0.16, fill=CORAL)
    text(slide, 1.0, 1.3, 11, 0.5, [[R(f"PARTE {part}", 16, CORAL, True)]])
    text(slide, 0.95, 1.8, 11.5, 1.2, [[R(ttl, 44, TEXT, True)]])
    rect(slide, 1.02, 2.85, 1.4, 0.05, fill=CORAL)
    text(slide, 1.0, 3.05, 11.3, 0.6, [[R(subtitle, 15, MUTED, False)]], line_spacing=1.1)
    y = 3.85
    for i, it in enumerate(items):
        r, c = divmod(i, 2)
        x = 1.0 + c * 6.0
        yy = y + r * 0.62
        rect(slide, x, yy, 5.7, 0.52, fill=PANEL, line=STROKE, rounded=True, radius=0.16)
        text(slide, x + 0.25, yy + 0.05, 5.3, 0.42, [[R(it, 12.5, TEXT, False)]],
             anchor=MSO_ANCHOR.MIDDLE)
    if page is not None:
        text(slide, 11.8, 7.02, 1.2, 0.35, [[R(str(page), 9.5, FAINT, True)]],
             align=PP_ALIGN.RIGHT)


# ----------------------------------------------------------------------------- build
def build():
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    blank = prs.slide_layouts[6]

    _pg = [0]

    def new():
        _pg[0] += 1
        return prs.slides.add_slide(blank), _pg[0]

    # =========================================================================
    # PORTADA + AGENDA
    # =========================================================================

    # ---- 1. Portada
    s, _ = new()
    rect(s, -0.1, -0.1, 13.53, 7.7, fill=BG)
    rect(s, 0, 0, 13.333, 0.16, fill=CORAL)
    rect(s, 0, 7.34, 13.333, 0.16, fill=CORAL)
    text(s, 1.0, 1.4, 11, 0.5, [[R("CURSO / WORKSHOP · DEVELOPER TOOLING", 14, CORAL, True)]])
    text(s, 0.95, 2.0, 11.5, 2.0,
         [[R("Claude Code", 60, TEXT, True)],
          [R("de asistente a ", 26, MUTED, False), R("sistema de desarrollo agéntico", 26, TEXT, True)]],
         space_after=10)
    rect(s, 1.02, 4.2, 1.4, 0.05, fill=CORAL)
    text(s, 1.0, 4.45, 11.5, 1.2,
         [[R("PARTE 1 — ", 15, CORAL, True),
           R("Claude Code: instalación · memoria · contexto & caching · MCP · skills · subagents & teams · automatización", 13.5, MUTED, False)],
          [R("PARTE 2 — ", 15, BLUE, True),
           R("La metodología (agnóstica): flujo con gates · CodeGraph / Serena / GSD · grafo de tickets · transferencia · ops", 13.5, MUTED, False)]],
         space_after=8)
    text(s, 1.0, 6.7, 11, 0.4,
         [[R("Basado en la documentación oficial (code.claude.com), el curso de hooks de Anthropic, "
             "GSD, CodeGraph y una instalación real en producción.", 11, FAINT, False)]])

    # ---- 2. Agenda
    s, pg = new()
    base(s, "Agenda", "Un curso, dos partes", page=pg)
    text(s, 0.7, 1.95, 5.85, 0.4, [[R("PARTE 1 · Claude Code", 14, CORAL, True)]])
    p1 = ["01  Instalación y uso básico", "02  Memoria, instrucciones y sesiones",
          "03  Contexto: context window y prompt caching", "04  MCP — conectar tus herramientas",
          "05  Plugins, tools y skills", "06  Subagents y Agent Teams", "07  Automatización"]
    for i, it in enumerate(p1):
        y = 2.4 + i * 0.56
        rect(s, 0.7, y, 5.85, 0.48, fill=PANEL, line=STROKE, rounded=True, radius=0.14)
        text(s, 0.95, y + 0.04, 5.5, 0.4, [[R(it, 12, TEXT, False)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 6.8, 1.95, 5.85, 0.4, [[R("PARTE 2 · La metodología (agnóstica)", 14, BLUE, True)]])
    p2 = ["08  El flujo de 11 etapas + ejemplo real", "09  Las herramientas del método",
          "10  El grafo de conocimiento de tickets", "11  Transferir la metodología (Copilot)",
          "12  Sincronización de máquinas", "13  Cierre"]
    for i, it in enumerate(p2):
        y = 2.4 + i * 0.56
        rect(s, 6.8, y, 5.85, 0.48, fill=PANEL, line=STROKE, rounded=True, radius=0.14)
        text(s, 7.05, y + 0.04, 5.5, 0.4, [[R(it, 12, TEXT, False)]], anchor=MSO_ANCHOR.MIDDLE)

    # =========================================================================
    # PARTE 1 — CLAUDE CODE
    # =========================================================================

    # ---- 3. Separador Parte 1
    s, pg = new()
    divider(s, 1, "Claude Code", "La herramienta: del binario en tu terminal a equipos de agentes.",
            ["01 · Instalación y uso básico", "02 · Memoria y sesiones",
             "03 · Contexto y prompt caching", "04 · MCP",
             "05 · Plugins, tools y skills", "06 · Subagents y Agent Teams",
             "07 · Automatización"], page=pg)

    # ---- 4. Instalación
    s, pg = new()
    base(s, "Parte 1 · 01 · Instalación y uso básico", "Instalar es una línea", page=pg)
    cards(s, [
        ("1", "Native (recomendado)", "curl -fsSL https://claude.ai/install.sh | bash — macOS/Linux/WSL, con auto-update."),
        ("2", "Homebrew", "brew install --cask claude-code (brew upgrade para actualizar)."),
        ("3", "Windows", "winget install Anthropic.ClaudeCode · o el instalador PowerShell."),
        ("4", "Linux pkg", "apt / dnf / apk en Debian, Fedora, RHEL, Alpine."),
    ], top=2.05, cols=2, height=1.0)
    code(s, 0.7, 4.35, 5.9, 1.35,
         ["cd tu-proyecto", "claude            # login la 1ª vez"],
         title="Arrancar en cualquier proyecto")
    panel_bullets(s, 6.9, 4.01, 5.75, 1.7, "El mismo motor en todas partes",
                  ["Terminal · VS Code / Cursor · JetBrains", "Desktop (diffs, sesiones en paralelo)",
                   "Web (claude.ai/code, tareas largas)"], accent=BLUE)
    takeaway(s, "Instalar es trivial; el valor está en cómo lo usas.")

    # ---- 5. Los dos modos
    s, pg = new()
    base(s, "Parte 1 · 01 · Instalación y uso básico", "Dos modos: interactivo y headless", page=pg)
    panel_bullets(s, 0.7, 2.05, 5.85, 2.5, "Interactivo — pensar contigo",
                  ["Abres `claude` y conversas",
                   "Plan mode: propone un plan y tú lo apruebas antes de tocar nada",
                   "Ideal para explorar, diseñar, depurar"], accent=CORAL)
    panel_bullets(s, 6.8, 2.05, 5.85, 2.5, "Headless (-p) — componible",
                  ["Un prompt: entra por stdin, sale por stdout",
                   "Encaja en tuberías Unix y en CI",
                   "La base de toda la automatización"], accent=BLUE)
    code(s, 0.7, 5.05, 11.95, 1.1,
         ["tail -200 app.log | claude -p \"avísame si ves anomalías\"",
          "git diff main --name-only | claude -p \"revisa estos ficheros por seguridad\""],
         title="El mismo Claude, dos formas de invocarlo")

    # ---- 6. Memoria: CLAUDE.md dos niveles
    s, pg = new()
    base(s, "Parte 1 · 02 · Memoria e instrucciones", "CLAUDE.md: el patrón de dos niveles", page=pg)
    panel_bullets(s, 0.7, 2.05, 5.85, 3.0, "Nivel 1 — siempre cargado (pequeño)",
                  ["Orientación mínima: qué es el proyecto, mapa de repos, comandos",
                   "Solo PUNTEROS de una línea a todo lo demás",
                   "Es la onboarding de 30s del agente",
                   "Regla write-once: el core apunta, no copia"], accent=CORAL)
    panel_bullets(s, 6.8, 2.05, 5.85, 3.0, "Nivel 2 — bajo demanda (el detalle)",
                  ["data/changes/STATUS.md — estado vivo por ticket",
                   "PLAYBOOK.md — lecciones de debugging",
                   "SHARP_EDGES.md — invariantes 'no tocar'",
                   "<TICKET>/<TICKET>.md — ~55 ledgers autocontenidos"], accent=BLUE)
    takeaway(s, "Recortamos el CLAUDE.md ~73% sin perder información: punteros, no copias.")

    # ---- 7. Jerarquía + settings + auto-memory
    s, pg = new()
    base(s, "Parte 1 · 02 · Instrucciones, permisos y auto-memory", "Jerarquía, permisos y memoria automática", page=pg)
    code(s, 0.7, 2.15, 6.0, 2.2,
         ["~/.claude/CLAUDE.md      # global usuario",
          "<repo>/CLAUDE.md         # compartido (versionado)",
          "<repo>/sub/CLAUDE.md     # al entrar en la subcarpeta",
          "<repo>/CLAUDE.local.md   # personal, no versionado"],
         title="Jerarquía de CLAUDE.md (se combinan; @fichero importa)")
    panel_bullets(s, 7.0, 1.82, 5.65, 2.55, "Permisos = allowlist específico",
                  ["settings.json: allow / deny / ask",
                   "Reglas concretas, no Bash(*)",
                   "mcp__serena__find_symbol, pytest exacto…",
                   "El humano es dueño de push/PR/deploy"], accent=GREEN)
    panel_bullets(s, 0.7, 4.55, 11.95, 1.55, "Auto-memory",
                  ["Claude guarda aprendizajes (comandos de build, pistas de debug) entre sesiones sin que escribas nada."],
                  accent=BLUE)

    # ---- 8. Sesiones
    s, pg = new()
    base(s, "Parte 1 · 02 · Sesiones", "Sesiones que viajan entre superficies", page=pg)
    cards(s, [
        ("→", "Teleport", "claude --teleport trae una sesión web/móvil al terminal."),
        ("↺", "Remote Control", "Sigue una sesión local desde el móvil o cualquier navegador."),
        ("▣", "Desktop handoff", "/desktop pasa la sesión al Desktop para revisar diffs."),
        ("⌁", "Channels / Slack", "Empuja eventos externos o convierte un bug de Slack en un PR."),
    ], top=2.15, cols=2, height=1.1)
    takeaway(s, "El mismo motor y los mismos CLAUDE.md/settings/MCP en todas las superficies.")

    # ---- 9. Contexto: anatomía + comandos
    s, pg = new()
    base(s, "Parte 1 · 03 · Contexto", "Context window: el recurso que gobierna todo", page=pg)
    rows = [
        ("System prompt", "~4.200 tokens — oculto, siempre primero", "fijo"),
        ("Auto-memory + entorno", "~680 + ~280 tokens al arrancar", "fijo"),
        ("Tools MCP", "índice ~120 tokens; el schema, al usar la tool", "por server"),
        ("Tu CLAUDE.md", "lo que tú decidas — por eso el patrón de dos niveles", "tú decides"),
        ("Conversación + tool results", "ficheros leídos, output de comandos… crece cada turno", "crece"),
    ]
    y0 = 2.0
    for i, (k, v, tag) in enumerate(rows):
        y = y0 + i * 0.58
        rect(s, 0.7, y, 7.6, 0.48, fill=PANEL, line=STROKE, rounded=True, radius=0.12)
        text(s, 0.95, y + 0.04, 2.6, 0.4, [[R(k, 11.5, BLUE, True)]], anchor=MSO_ANCHOR.MIDDLE)
        text(s, 3.6, y + 0.04, 4.0, 0.4, [[R(v, 10.5, TEXT, False)]], anchor=MSO_ANCHOR.MIDDLE)
        text(s, 7.35, y + 0.04, 0.85, 0.4, [[R(tag, 9, MUTED, False)]],
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    code(s, 8.6, 2.35, 4.05, 2.55,
         ["/context   # desglose de uso",
          "/compact céntrate en la API",
          "/clear     # reset entre tareas",
          "/rewind    # checkpoints (Esc+Esc)"],
         title="Los mandos")
    text(s, 0.7, 5.05, 11.95, 0.9,
         [[R("Ventana: 200K tokens (1M en beta vía API). ", 12, TEXT, True),
           R("El rendimiento degrada ANTES de llenarla: contexto con ruido = peores decisiones. "
             "Auto-compact salta cerca del límite y es lossy — mejor /compact manual con foco.",
             12, MUTED, False)]], line_spacing=1.1)
    takeaway(s, "Mide con /context antes de optimizar: te dice exactamente qué bloque engorda.")

    # ---- 10. Contexto: higiene
    s, pg = new()
    base(s, "Parte 1 · 03 · Contexto", "Higiene de contexto: máxima señal por token", page=pg)
    cards(s, [
        ("1", "CLAUDE.md mínimo", "Regla oficial: si puedes borrarlo sin que Claude se equivoque, bórralo. Dos niveles + punteros."),
        ("2", "Carga bajo demanda", "@imports y CLAUDE.md de subcarpeta solo cargan cuando tocan. El detalle vive fuera."),
        ("3", "MCP con moderación", "Cada server suma su bloque de tools. Desactiva los que el proyecto no use."),
        ("4", "Subagentes para investigar", "La exploración sucia va a un subagente (sección 06); vuelve solo el resumen."),
        ("5", "Explore → Plan → Code", "Separa exploración (plan mode) de implementación; el ruido no se queda a vivir."),
        ("6", "Lecturas con puntería", "“lee config/auth.js” > “entiende el auth”. Adelanto: CodeGraph-primero (Parte 2)."),
    ], top=2.0, cols=2, height=1.18)
    takeaway(s, "La prevalencia de tools de la Parte 2 es, en el fondo, política de contexto.")

    # ---- 11. Prompt caching: mecanismo
    s, pg = new()
    base(s, "Parte 1 · 03 · Contexto", "Prompt caching: no pagar lo mismo dos veces", page=pg)
    text(s, 0.7, 1.9, 11.95, 0.55,
         [[R("Cada turno reenvía TODO el contexto. La API cachea el prefijo estable — orden estricto ",
             12.5, MUTED, False),
           R("Tools → System → Messages", 12.5, CORAL, True, MONO),
           R(": un cambio invalida su nivel y los siguientes.", 12.5, MUTED, False)]])
    # tabla de precios
    prices = [
        ("TTL 5 min (defecto)", "escribir 1.25×", "leer 0.1×"),
        ("TTL 1 hora", "escribir 2×", "leer 0.1×"),
    ]
    for i, (a, b, c) in enumerate(prices):
        y = 2.6 + i * 0.6
        rect(s, 0.7, y, 5.85, 0.5, fill=PANEL, line=STROKE, rounded=True, radius=0.12)
        text(s, 0.95, y + 0.05, 2.5, 0.4, [[R(a, 11.5, BLUE, True)]], anchor=MSO_ANCHOR.MIDDLE)
        text(s, 3.5, y + 0.05, 1.7, 0.4, [[R(b, 11.5, TEXT, False)]], anchor=MSO_ANCHOR.MIDDLE)
        text(s, 5.2, y + 0.05, 1.2, 0.4, [[R(c, 11.5, GREEN, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 0.7, 3.95, 5.85, 1.6,
         [[R("Una sesión de 50 turnos relee el prefijo 50 veces a 0.1× — ", 12, TEXT, False),
           R("ese es el descuento que hace viables las sesiones largas.", 12, GREEN, True)],
          [R("Mínimo cacheable ~1k tokens · máx. 4 breakpoints · diagnóstico en usage.cache_read_input_tokens.",
             10.5, MUTED, False)]], space_after=8, line_spacing=1.1)
    code(s, 6.8, 2.6, 5.85, 2.9,
         ["system=[{ \"type\": \"text\",",
          "  \"text\": STABLE_INSTRUCTIONS,   # estable",
          "  \"cache_control\":",
          "    {\"type\": \"ephemeral\"} }]   # breakpoint",
          "messages=[{\"role\": \"user\",",
          "  \"content\": query }]  # variable: FUERA",
          "# error clásico: timestamp ANTES del",
          "# breakpoint -> 0 hits y nadie sabe por qué"],
         title="Lo estable primero; el breakpoint al final de lo estable")
    takeaway(s, "Demo ejecutable: ejemplos/prompt-caching/cache_demo.py — mira los contadores de cache.")

    # ---- 12. Prompt caching en Claude Code
    s, pg = new()
    base(s, "Parte 1 · 03 · Contexto", "Caching en Claude Code: tú decides el hit-rate", page=pg)
    rows = [
        ("CLAUDE.md pequeño y estable", "Prefijo corto que nunca cambia → hits en cada turno", True),
        ("Editar CLAUDE.md/settings a mitad de sesión", "Cambia el prefijo → invalida el cache desde ahí", False),
        ("Muchos servers MCP activos", "Bloque de tools grande y cambiante → escrituras caras", False),
        ("Sesión larga, turnos frecuentes", "El historial se relee a 0.1× — el caso ideal", True),
        ("/compact o auto-compact", "Reescribe el historial → rompe el cache de mensajes una vez, y sigue", False),
    ]
    y0 = 2.05
    for i, (k, v, good) in enumerate(rows):
        y = y0 + i * 0.68
        acc = GREEN if good else CORAL
        rect(s, 0.7, y, 11.95, 0.58, fill=PANEL, line=STROKE, rounded=True, radius=0.1)
        rect(s, 0.7, y, 0.08, 0.58, fill=acc, rounded=True, radius=0.5)
        text(s, 0.95, y + 0.05, 4.6, 0.48, [[R(k, 11.5, acc, True)]], anchor=MSO_ANCHOR.MIDDLE)
        text(s, 5.7, y + 0.05, 6.8, 0.48, [[R(v, 11.5, TEXT, False)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 0.7, 5.6, 11.95, 0.55,
         [[R("Claude Code no expone cache_control: lo aplica solo (system + tools + historial = prefijo). "
             "Tus decisiones de la slide anterior determinan el hit-rate.", 11.5, MUTED, False)]])
    takeaway(s, "Contexto lean y estable rinde mejor Y cuesta menos: dos niveles optimiza ambos ejes a la vez.")

    # ---- 13. MCP
    s, pg = new()
    base(s, "Parte 1 · 04 · MCP", "Model Context Protocol: conecta tu mundo", page=pg)
    text(s, 0.7, 1.9, 11.95, 0.6,
         [[R("Estándar abierto para enchufar Claude a datos y tools externas. Las tools aparecen como ",
             12.5, MUTED, False), R("mcp__<server>__<tool>", 12.5, CORAL, True, MONO), R(".", 12.5, MUTED, False)]])
    cards(s, [
        ("L", "Scope local", ".claude/settings.local.json — solo tú, esta máquina."),
        ("P", "Scope project", ".mcp.json en la raíz, versionado, compartido con el equipo."),
        ("U", "Scope user", "~/.claude.json — todos tus proyectos."),
        ("★", "Servers que uso", "serena · context7 · playwright · codegraph · supabase."),
    ], top=2.55, cols=2, height=1.0)
    code(s, 0.7, 4.95, 11.95, 1.15,
         ["claude mcp add --transport http context7 https://mcp.context7.com/mcp",
          "claude mcp list      # estado      /mcp   # auth OAuth y tools dentro de la sesión"],
         title="Añadir un server (secretos por env var; cada server suma contexto — sección 03)")

    # ---- 14. Plugins, tools y skills
    s, pg = new()
    base(s, "Parte 1 · 05 · Plugins, tools y skills", "Cuatro capas de extensibilidad", page=pg)
    cards(s, [
        ("1", "Tools", "Lo que Claude puede hacer: Read/Edit/Bash/Grep/Task + mcp__*. Gobernadas por permisos."),
        ("2", "Slash commands", ".claude/commands/x.md = /x. El cuerpo es el prompt. Lo invocas tú."),
        ("3", "Skills", ".claude/skills/x/SKILL.md con description. Claude la auto-selecciona; lleva scripts/plantillas."),
        ("4", "Plugins + marketplace", "Agrupan skills+comandos+agentes+MCP+hooks. /plugin install. GSD es uno."),
    ], top=2.05, cols=2, height=1.15)
    text(s, 0.7, 4.75, 11.95, 0.9,
         [[R("La 5ª capa — subagentes y equipos de agentes — merece sección propia: es la siguiente.",
             12.5, BLUE, True)]])
    takeaway(s, "Slash command = atajo que invocas tú. Skill = capacidad que Claude decide usar. Plugin = el reparto.")

    # ---- 15. Subagents
    s, pg = new()
    base(s, "Parte 1 · 06 · Subagents y Agent Teams", "Subagentes: el ruido muere fuera de tu sesión", page=pg)
    panel_bullets(s, 0.7, 2.0, 5.85, 2.2, "Built-ins (tool Task)",
                  ["Explore — búsqueda read-only por el codebase",
                   "Plan — diseñar la estrategia",
                   "general-purpose — multi-paso, todos los tools",
                   "En paralelo para trabajo independiente"], accent=BLUE)
    code(s, 6.8, 2.35, 5.85, 2.6,
         ["# .claude/agents/security-reviewer.md",
          "---",
          "name: security-reviewer",
          "description: Revisa vulnerabilidades",
          "  tras cambios en auth o deps.",
          "tools: Read, Grep, Glob, Bash",
          "model: opus",
          "---",
          "Eres un ingeniero de seguridad senior…"],
         title="Custom: un .md con frontmatter (/agents los lista)")
    panel_bullets(s, 0.7, 4.35, 5.85, 1.7, "Lo que hay que saber",
                  ["Contexto AISLADO: solo el resumen vuelve a tu sesión",
                   "NO hereda tu conversación: contexto en el prompt",
                   "description = auto-selección · tools = allowlist propio"], accent=CORAL)
    takeaway(s, "Ejemplos reales en ejemplos/subagents/: security-reviewer y refactor-scout (regla CodeGraph→Serena).")

    # ---- 16. Agent Teams
    s, pg = new()
    base(s, "Parte 1 · 06 · Subagents y Agent Teams", "Agent Teams: varios Claudes que se coordinan", page=pg)
    panel_bullets(s, 0.7, 2.0, 5.85, 2.6, "La arquitectura (experimental)",
                  ["Lead + teammates: cada uno una SESIÓN completa",
                   "Task list compartida: ~/.claude/tasks/<team>/",
                   "Mensajería directa entre teammates (inboxes)",
                   "Pueden debatir hallazgos, no solo reportar",
                   "Roles = tus subagentes custom reutilizados"], accent=CORAL)
    code(s, 6.8, 2.35, 5.85, 1.6,
         ["// ~/.claude/settings.json",
          "{ \"env\": { \"CLAUDE_CODE_EXPERIMENTAL",
          "_AGENT_TEAMS\": \"1\" },",
          "  \"teammateMode\": \"in-process\" }  // tmux|iterm2"],
         title="Activar (off por defecto)")
    panel_bullets(s, 6.8, 4.25, 5.85, 1.8, "Límites hoy",
                  ["Sin /resume in-process · un team por sesión · sin anidar",
                   "Split-panes requiere tmux o iTerm2",
                   "Coste: cada teammate es una sesión completa"], accent=BLUE)
    text(s, 0.7, 4.75, 5.85, 1.3,
         [[R("Se pide en lenguaje natural: ", 11.5, MUTED, False),
           R("“monta un equipo con un architect y dos implementers; exige aprobación de plan antes de tocar código”",
             11.5, TEXT, False)]], line_spacing=1.1)
    takeaway(s, "Particiona los ficheros: cada teammate es dueño de los suyos.")

    # ---- 17. Subagente vs team
    s, pg = new()
    base(s, "Parte 1 · 06 · Subagents y Agent Teams", "¿Subagente o team? La decisión", page=pg)
    hdr_y = 2.0
    text(s, 4.6, hdr_y, 3.7, 0.4, [[R("SUBAGENTE", 13, BLUE, True)]], align=PP_ALIGN.CENTER)
    text(s, 8.6, hdr_y, 3.7, 0.4, [[R("AGENT TEAM", 13, CORAL, True)]], align=PP_ALIGN.CENTER)
    rows = [
        ("Contexto", "Aislado; devuelve un resumen", "Cada teammate = sesión completa"),
        ("Comunicación", "Solo resultado → principal", "Task list + mensajes entre teammates"),
        ("Coste", "Bajo (lo caro muere fuera)", "Alto (N sesiones en paralelo)"),
        ("Úsalo para", "Side-quests: investigar, verificar", "Paralelismo real: revisión multi-capa"),
        ("Config", "Nada, o .claude/agents/*.md", "Env var experimental + settings"),
    ]
    for i, (k, a, b) in enumerate(rows):
        y = 2.45 + i * 0.62
        rect(s, 0.7, y, 11.95, 0.52, fill=PANEL, line=STROKE, rounded=True, radius=0.1)
        text(s, 0.95, y + 0.05, 3.3, 0.42, [[R(k, 11.5, TEXT, True)]], anchor=MSO_ANCHOR.MIDDLE)
        text(s, 4.6, y + 0.05, 3.7, 0.42, [[R(a, 11, MUTED, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, 8.6, y + 0.05, 3.7, 0.42, [[R(b, 11, MUTED, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 0.7, 5.68, 11.95, 0.5,
         [[R("Puente a la Parte 2: ", 11.5, GREEN, True),
           R("GSD es esto en producción — gsd-planner, gsd-executor, gsd-verifier son subagentes custom "
             "distribuidos como plugin.", 11.5, MUTED, False)]])
    takeaway(s, "Subagente para que el ruido muera fuera; team para que varios Claudes debatan. El coste no es el mismo.")

    # ---- 18. Hooks
    s, pg = new()
    base(s, "Parte 1 · 07 · Automatización", "Hooks: el control determinista", page=pg)
    panel_bullets(s, 0.7, 2.0, 5.85, 3.1, "El contrato de un hook",
                  ["Payload del evento por STDIN",
                   "Eventos: PreToolUse, PostToolUse, SessionStart, Stop…",
                   "Matcher = regex sobre el nombre de la tool",
                   "exit 0 permite · exit 2 BLOQUEA (stderr → Claude)",
                   "Pre = intención (tool_input) · Post = resultado (tool_response)"], accent=CORAL)
    code(s, 6.8, 2.35, 5.85, 2.4,
         ["if (payload.tool_input",
          "     ?.file_path?.includes(\".env\")) {",
          "  console.error(\"Bloqueado: no leas .env\");",
          "  process.exit(2);   // <-- veto",
          "}",
          "process.exit(0);"],
         title="Patrón de bloqueo (JS)")
    takeaway(s, "Con instrucciones le pides que se porte bien; con un hook lo garantizas.")

    # ---- 19. Automatización: niveles
    s, pg = new()
    base(s, "Parte 1 · 07 · Automatización", "De hooks a agentes programados", page=pg)
    cards(s, [
        ("a", "Hooks", "Seguridad (.env), formato (prettier), type-check bloqueante, 'IA revisando IA' vía SDK."),
        ("b", "Headless / piping", "claude -p en cualquier tubería Unix o script."),
        ("c", "CI/CD", "GitHub Actions · GitLab · GitHub Code Review sobre cada PR."),
        ("d", "Scheduling", "Routines (/schedule, infra Anthropic) · Desktop tasks · /loop."),
    ], top=2.05, cols=2, height=1.05)
    code(s, 0.7, 4.5, 11.95, 1.25,
         ["import { query } from \"@anthropic-ai/claude-agent-sdk\";",
          "for await (const m of query({ prompt, options: { allowedTools: [\"Edit\"] } }))",
          "  if (m.type === \"result\") console.log(m.result);"],
         title="Agent SDK — workflows a medida (menor privilegio)")

    # =========================================================================
    # PARTE 2 — LA METODOLOGÍA
    # =========================================================================

    # ---- 20. Separador Parte 2
    s, pg = new()
    divider(s, 2, "La metodología",
            "Agnóstica de la herramienta: se demuestra con Claude Code, pero la disciplina viaja.",
            ["08 · El flujo de 11 etapas + ejemplo real", "09 · Las herramientas del método",
             "10 · El grafo de conocimiento de tickets", "11 · Transferir a otro agente (Copilot)",
             "12 · Sincronización de máquinas", "13 · Cierre"], page=pg)

    # ---- 21. Principio
    s, pg = new()
    base(s, "Parte 2 · 08 · La metodología", "El agente es un colaborador disciplinado", page=pg)
    rect(s, 0.7, 2.0, 11.95, 1.15, fill=PANEL_2, line=CORAL, lw=1.0, rounded=True, radius=0.06)
    text(s, 1.0, 2.18, 11.4, 0.85,
         [[R("“La autonomía se gana por-decisión, no se concede en bloque.”", 17, TEXT, True)]],
         anchor=MSO_ANCHOR.MIDDLE)
    panel_bullets(s, 0.7, 3.4, 5.85, 2.05, "El agente posee",
                  ["Investigación y diagnóstico", "Planes e implementación", "Tests y documentación"], accent=BLUE)
    panel_bullets(s, 6.8, 3.4, 5.85, 2.05, "El humano posee",
                  ["Decisiones go/no-go y el scope", "Toda acción externa: push, PR, deploy",
                   "La aprobación del plan"], accent=GREEN)
    text(s, 0.72, 5.62, 11.9, 0.5,
         [[R("Agnóstica:  ", 11.5, CORAL, True),
           R("no es un flujo de Claude Code — es una forma de trabajar con CUALQUIER agente. "
             "La sección 11 lo demuestra transfiriéndola a GitHub Copilot.", 11, MUTED, False)]],
         line_spacing=1.03)
    takeaway(s, "El modelo no gana confianza gratis: la gana decisión a decisión, con evidencia.")

    # ---- 22. El flujo de 11 etapas
    s, pg = new()
    base(s, "Parte 2 · 08 · La metodología", "El flujo real de 11 etapas", page=pg)
    stages = [
        ("1", "Orientar: /kg + history-first + status-first", False),
        ("2", "Triaje inbound: ¿síntoma en el contrato?", True),
        ("3", "Regresión vs. pre-existente (probarlo)", True),
        ("4", "Investigar: oráculo determinista barato", False),
        ("5", "Plan → acuerdo humano explícito", True),
        ("6", "Implementar: TDD RED → GREEN, mínimo", False),
        ("7", "Verificar: contrato (wrapper) + imagen desplegada", True),
        ("8", "Documentar — cada cosa una vez", False),
        ("9", "Sanitizar — líneas añadidas", False),
        ("10", "Handoff: el humano hace push / PR", False),
        ("11", "Revisión bot + persistir lecciones", True),
    ]
    for i, (n, label, gate) in enumerate(stages):
        col = i // 6
        row = i % 6
        x = 0.7 + col * 6.15
        y = 2.05 + row * 0.66
        acc = CORAL if gate else PANEL_2
        rect(s, x, y, 5.85, 0.56, fill=PANEL, line=(CORAL if gate else STROKE),
             lw=(1.0 if gate else 0.75), rounded=True, radius=0.14)
        rect(s, x + 0.14, y + 0.11, 0.34, 0.34, fill=acc, rounded=True, radius=0.4)
        text(s, x + 0.14, y + 0.11, 0.34, 0.34,
             [[R(n, 11, BG if gate else TEXT, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, x + 0.62, y + 0.06, 5.1, 0.46, [[R(label, 11.5, TEXT, False)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 0.7, 6.02, 11.95, 0.3,
         [[R("■ ", 11, CORAL, True), R("Los recuadros coral son GATES (puntos de decisión). "
           "Un gate rojo = STOP: no escribir código.", 11, MUTED, False)]])
    takeaway(s, "Ninguna etapa dice 'pídele al LLM que lo arregle': la potencia se canaliza por gates.")

    # ---- 23. Ejemplo real
    s, pg = new()
    base(s, "Parte 2 · 08 · La metodología", "Un ejemplo real, de principio a fin", page=pg)
    rect(s, 0.7, 1.95, 11.95, 0.62, fill=PANEL_2, line=CORAL, rounded=True, radius=0.1)
    text(s, 0.95, 2.0, 11.5, 0.52,
         [[R("Bug QA: ", 12.5, CORAL, True),
           R("“un campo sale vacío en la UI, pero está escrito en el PDF.”", 12.5, TEXT, False)]],
         anchor=MSO_ANCHOR.MIDDLE)
    steps = [
        ("Orientar", "/kg saca la zona de peligro (un SHARP_EDGE que restringe el fix); git/gh dan la base."),
        ("Contrato", "El campo está vacío en el JSON del endpoint (Playwright/F12) → el Lambda es responsable."),
        ("Provenance", "Re-extraer en la base previa: ya salía vacío → pre-existente, no regresión."),
        ("Investigar", "Serena + CodeGraph localizan el detector; un _diag_pdf.py determinista da la causa: SIN LLM."),
        ("Fix + verificar", "Test RED → fix estructural (no cliente) → regresión no-op + contrato local (wrapper) + dentro de la imagen desplegada."),
        ("Handoff", "Sanitizar → el humano hace push/PR. Un bot detecta un caso más → test + PLAYBOOK."),
    ]
    for i, (tag, desc) in enumerate(steps):
        y = 2.78 + i * 0.56
        rect(s, 0.7, y, 11.95, 0.48, fill=PANEL, line=STROKE, rounded=True, radius=0.12)
        rect(s, 0.7, y, 0.08, 0.48, fill=CORAL, rounded=True, radius=0.5)
        text(s, 0.95, y + 0.04, 2.4, 0.4, [[R(tag, 11.5, BLUE, True)]], anchor=MSO_ANCHOR.MIDDLE)
        text(s, 3.5, y + 0.04, 9.0, 0.4, [[R(desc, 11, TEXT, False)]], anchor=MSO_ANCHOR.MIDDLE)
    takeaway(s, "El modelo no gana confianza gratis: la gana decisión a decisión, con evidencia.")

    # ---- 24. Prevalencia de tools
    s, pg = new()
    base(s, "Parte 2 · 09 · Herramientas del método", "Qué herramienta usa el agente, y cuándo", page=pg)
    text(s, 0.7, 1.9, 11.95, 0.5,
         [[R("Regla del CLAUDE.md: ", 12, MUTED, False),
           R("codegraph_explore PRIMERO (survey en 1 llamada) · Serena find_referencing_symbols "
             "para el chequeo preciso antes de renombrar. ", 12, TEXT, True),
           R("Orden: barato → caro, determinista → probabilístico.", 12, MUTED, False)]])
    rows = [
        ("Orientar", "/kg (grafo de tickets → zona de peligro) · STATUS.md · git · gh", "coste 0"),
        ("Navegar (survey)", "CodeGraph codegraph_explore — fuente + rutas + blast radius + cobertura", "1 llamada"),
        ("Refactor-check", "Serena find_referencing_symbols — desambigua por clase (antes de renombrar)", "preciso"),
        ("Diagnosticar", "Oráculo determinista: parser / validador / _diag_*.py", "coste 0"),
        ("Entorno (logs, config)", "AWS CLI — CloudWatch · lambda get-function · SQS/DLQ", "read-only"),
        ("Contrato de salida", "Playwright / F12 sobre el endpoint que ve el consumidor", "reproducir"),
        ("Solo al final", "La tirada del LLM — para VERIFICAR el fix, no para diagnosticar", "metered"),
    ]
    y0 = 2.5
    for i, (etapa, tool, tag) in enumerate(rows):
        y = y0 + i * 0.6
        last = (i == len(rows) - 1)
        rect(s, 0.7, y, 11.95, 0.5, fill=(PANEL_2 if last else PANEL),
             line=(CORAL if last else STROKE), lw=(1.0 if last else 0.75), rounded=True, radius=0.12)
        text(s, 0.95, y + 0.05, 2.9, 0.4, [[R(etapa, 11.5, (CORAL if last else BLUE), True)]],
             anchor=MSO_ANCHOR.MIDDLE)
        text(s, 3.95, y + 0.05, 7.0, 0.4, [[R(tool, 11.5, TEXT, False)]], anchor=MSO_ANCHOR.MIDDLE)
        text(s, 11.1, y + 0.05, 1.4, 0.4, [[R(tag, 9.5, MUTED, False)]],
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    # ---- 25. Qué es cada herramienta → mapa a fases
    s, pg = new()
    base(s, "Parte 2 · 09 · Herramientas del método", "Cada herramienta, en una frase — y su fase", page=pg)
    cards(s, [
        ("CG", "CodeGraph → investigar", "Grafo del CÓDIGO (tree-sitter→SQLite, local). Survey en 1 llamada: fuente + callers + blast radius + cobertura."),
        ("Se", "Serena → pre-refactor", "Navegación semántica vía LSP. find_referencing_symbols desambigua por clase: el chequeo preciso antes de renombrar."),
        ("G", "GSD → plan/ejecutar/verificar", "El método hecho tooling: subagentes gsd-* + estado en .planning/. Automatiza los gates 5-7."),
        ("/kg", "Grafo de tickets → orientar", "Grafo de la MEMORIA del proyecto: tickets + sharp edges. History-first sin LLM (sección 10)."),
        ("Pw", "Playwright → triaje y outbound", "Reproduce el síntoma donde lo ve el consumidor: el endpoint real, no una función interna."),
        ("Or", "Oráculos → diagnosticar", "Parsers/validadores/_diag_*.py propios: respuesta barata y reproducible antes de gastar la tirada del LLM."),
    ], top=2.0, cols=2, height=1.18)
    takeaway(s, "La inversión clásica — el modelo diagnostica — es lo que este orden evita: el modelo verifica; los oráculos diagnostican.")

    # ---- 26. GSD
    s, pg = new()
    base(s, "Parte 2 · 09 · Herramientas del método", "GSD: la metodología hecha tooling", page=pg)
    code(s, 0.7, 2.2, 6.15, 3.2,
         ["/gsd-new-project  -> PROJECT.md + ROADMAP.md",
          "/gsd-plan-phase   -> PLAN.md  (+ gate",
          "                     gsd-plan-checker)",
          "/gsd-execute-phase-> olas paralelas,",
          "                     commits atómicos",
          "/gsd-verify-work  -> VERIFICATION.md",
          "",
          "/gsd-progress     # qué toca ahora"],
         title="Ciclo por fases, estado versionado en .planning/")
    panel_bullets(s, 7.1, 1.86, 5.55, 3.55, "Subagentes especializados",
                  ["gsd-planner — desglose + dependencias",
                   "gsd-plan-checker — ¿logrará el objetivo? (goal-backward)",
                   "gsd-executor — commits atómicos, checkpoints",
                   "gsd-code-reviewer — REVIEW.md por severidad",
                   "gsd-verifier — verifica el OBJETIVO, no solo tareas"], accent=CORAL)
    text(s, 0.7, 5.55, 11.95, 0.5,
         [[R("Es la Parte 1 aplicada: ", 11.5, GREEN, True),
           R("subagentes custom (sección 06) + skills, distribuidos como plugin (sección 05).",
             11.5, MUTED, False)]])
    takeaway(s, "GSD automatiza los gates que la metodología aplica a mano: plan, verificación, rastro durable.")

    # ---- 27. CodeGraph
    s, pg = new()
    base(s, "Parte 2 · 09 · Herramientas del método", "CodeGraph: inteligencia de código local", page=pg)
    panel_bullets(s, 0.7, 2.0, 5.85, 2.4, "Índice tree-sitter → SQLite (.codegraph/)",
                  ["Símbolos: funciones, clases, rutas, componentes",
                   "Aristas: llamadas, imports, herencia, referencias",
                   "Determinista (del AST), sin API keys",
                   "Devuelve: fuente + rutas + blast radius + cobertura"], accent=BLUE)
    code(s, 6.8, 2.2, 5.85, 1.7,
         ["codegraph init      # crea .codegraph/",
          "codegraph sync      # incremental tras editar",
          "codegraph explore \"<símbolo|pregunta>\"",
          "codegraph serve --path <repo> --mcp"],
         title="CLI  (install --location=global escribe el MCP)")
    rect(s, 6.8, 4.15, 5.85, 1.2, fill=PANEL, line=STROKE, rounded=True, radius=0.06)
    text(s, 7.05, 4.28, 5.4, 1.05,
         [[R("Primero para navegar:  ", 11, CORAL, True),
           R("fuente + callers + blast radius + cobertura en 1 consulta (trátala como YA leída). "
             "Serena find_referencing_symbols = chequeo preciso (desambigua por clase). "
             "--path fija el repo por defecto.", 11, TEXT, False)]], line_spacing=1.03)
    text(s, 0.7, 4.6, 5.85, 1.4,
         [[R("58% menos tool calls · 22% más rápido", 13, GREEN, True)],
          [R("(según sus benchmarks: casi elimina las lecturas de fichero)", 10.5, MUTED, False)]],
         space_after=4)
    takeaway(s, "Una consulta en vez de grep→abrir→seguir-import→repetir. Menos contexto, más señal.")

    # ---- 28. Serena
    s, pg = new()
    base(s, "Parte 2 · 09 · Herramientas del método", "Serena: navegación semántica, chequeo preciso", page=pg)
    panel_bullets(s, 0.7, 2.0, 5.85, 2.6, "Navegación semántica vía LSP (MCP)",
                  ["Opera sobre SÍMBOLOS, no texto: precisión de IDE",
                   "find_symbol (body=true) — un método de un fichero de 5k líneas",
                   "get_symbols_overview — esqueleto de un fichero",
                   "find_referencing_symbols — quién referencia, POR CLASE",
                   "search_for_pattern · activate_project (multi-repo)"], accent=BLUE)
    code(s, 6.8, 2.35, 5.85, 1.35,
         ["claude mcp add serena -- uvx \\",
          "  --from git+https://github.com/oraios/serena \\",
          "  serena start-mcp-server"],
         title="Instalar (stdio)")
    rect(s, 6.8, 4.0, 5.85, 1.55, fill=PANEL, line=CORAL, rounded=True, radius=0.06)
    text(s, 7.05, 4.15, 5.4, 1.3,
         [[R("Por qué es OBLIGATORIO pre-rename:  ", 11, CORAL, True),
           R("el impact plano de CodeGraph mezcla métodos homónimos (Invoice.process vs Refund.process); "
             "Serena los desambigua por clase. Complementarios, no rivales.", 11, TEXT, False)]],
         line_spacing=1.05)
    text(s, 0.7, 4.85, 5.85, 0.9,
         [[R("El subagente refactor-scout ", 11.5, GREEN, True),
           R("(ejemplos/subagents/) empaqueta el orden CodeGraph → Serena → grep como procedimiento.",
             11.5, MUTED, False)]], line_spacing=1.05)
    takeaway(s, "CodeGraph responde '¿qué se rompe?'; Serena responde '¿exactamente quién llama a ESTE process()?'")

    # ---- 29. KG: qué es y cómo se usa
    s, pg = new()
    base(s, "Parte 2 · 10 · Grafo de tickets", "CodeGraph, pero para tickets y lecciones", page=pg)
    panel_bullets(s, 0.7, 2.0, 5.85, 2.4, "Grafo sobre los .md del proyecto (data/, interno)",
                  ["Nodos: tickets + símbolos + invariantes + lecciones",
                   "Aristas tipadas: references / supersedes / relacionado",
                   "Comunidades = zonas de peligro (fin-de-carta, títulos…)",
                   "Determinista, sin LLM en la consulta (lee graph.json)"], accent=BLUE)
    code(s, 6.8, 2.2, 5.85, 1.9,
         ["/kg <ticket|tema>   # vecinos  (explain)",
          "/kg <A> <B>         # camino    (path)",
          "/kg find <substr>   # nombre de nodo",
          "/kg-refresh         # reconstruir el grafo"],
         title="Skill /kg  (envuelve graphify explain/path)")
    rect(s, 6.8, 4.28, 5.85, 1.12, fill=PANEL, line=STROKE, rounded=True, radius=0.06)
    text(s, 7.05, 4.4, 5.4, 0.95,
         [[R("Primero la historia:  ", 11, CORAL, True),
           R("corre /kg <ticket|tema> ANTES de grep — saca los tickets relacionados + la zona de "
             "peligro a leer. Apunta a qué leer, no lo sustituye.", 11, TEXT, False)]], line_spacing=1.03)
    text(s, 0.7, 4.6, 5.85, 1.4,
         [[R("Ejemplo real: /kg get_letter_end", 13, GREEN, True)],
          [R("devuelve al instante los 5-6 tickets que comparten ese código — la zona de peligro completa.",
             10.5, MUTED, False)]],
         space_after=4)
    takeaway(s, "El grafo es el mapa; el agente, el guía. Recall de zonas de peligro que si no tendrías que recordar.")

    # ---- 30. KG: cómo se construye y dónde se engancha
    s, pg = new()
    base(s, "Parte 2 · 10 · Grafo de tickets", "Cómo se construye — y dónde se engancha", page=pg)
    steps = [
        ("1", "Corpus curado", "manifest.txt DIFFEABLE: writeups sst-*, hubs (STATUS, SHARP_EDGES), runbooks, memoria. Sin binarios ni copias stale."),
        ("2", "prepare", "kg_refresh.sh: manifest → stage → copiar a un SCRATCH FUERA del repo (graphify respeta .gitignore y data/ lo está)."),
        ("3", "/graphify", "Extracción semántica con ~5 subagentes en paralelo → clustering. El único paso con LLM: en el BUILD, nunca en la consulta."),
        ("4", "finalize", "Copiar artefactos (graph.json tipado + graph.html + GRAPH_REPORT.md) + leak-check: nada fuera de data/."),
    ]
    for i, (n, t, d) in enumerate(steps):
        y = 2.0 + i * 0.82
        rect(s, 0.7, y, 11.95, 0.72, fill=PANEL, line=STROKE, rounded=True, radius=0.09)
        rect(s, 0.86, y + 0.14, 0.44, 0.44, fill=CORAL, rounded=True, radius=0.35)
        text(s, 0.86, y + 0.15, 0.44, 0.42, [[R(n, 13, BG, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, 1.5, y + 0.07, 1.9, 0.58, [[R(t, 12.5, BLUE, True)]], anchor=MSO_ANCHOR.MIDDLE)
        text(s, 3.5, y + 0.07, 9.0, 0.6, [[R(d, 10.5, TEXT, False)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    text(s, 0.7, 5.4, 11.95, 0.8,
         [[R("Dónde se engancha: ", 12, GREEN, True),
           R("etapa 1 (Orientar) — la regla history-first del CLAUDE.md. Honestidad: la ganancia es recall en "
             "zonas densas; EXTRACTED = fiable, INFERRED = pista. Artefacto derivado: nunca viaja, se reconstruye (sección 12).",
             11.5, MUTED, False)]], line_spacing=1.1)
    takeaway(s, "Un paso semántico en el build; cero LLM en la consulta. Por eso es un skill, no un script.")

    # ---- 31. Transferir la metodología
    s, pg = new()
    base(s, "Parte 2 · 11 · Transferencia", "La prueba de que es agnóstica: llevarla a Copilot", page=pg)
    panel_bullets(s, 0.7, 2.0, 5.85, 2.75, "Viaja SIN cambios (las 5 reglas)",
                  ["Plan → acuerdo → implementar",
                   "Verificar en el contrato del CONSUMIDOR",
                   "Resolver la clase general, no un input",
                   "Rastro durable: porqué, qué, cómo se verificó",
                   "El humano posee lo externo (merge, deploy)"], accent=GREEN)
    panel_bullets(s, 6.8, 2.0, 5.85, 2.75, "Se re-mapea por repo",
                  ["El contrato: HTTP / DB / evento / artefacto",
                   "El tracker: Jira · Azure Boards · Issues",
                   "La pirámide de tests y el runtime desplegado",
                   "Las reglas de sanitización locales",
                   "Las tools: Serena/CodeGraph//kg son fungibles"], accent=BLUE)
    panel_bullets(s, 0.7, 4.95, 11.95, 1.25, "El starter-kit (docs/ai-agents-code-methodology/)",
                  ["Plantillas STATUS · SHARP_EDGES · handover · QA + bootstrap-new-repo.ps1 → estructura en el repo destino",
                   "Sin grafo de tickets: fallback determinista (STATUS newest-first + búsqueda léxica + historia de commits) = 80% del valor"],
                  accent=CORAL)
    takeaway(s, "Si solo conservas cinco reglas, conserva esas cinco. Las tools se sustituyen; la disciplina viaja.")

    # ---- 32. Machine sync
    s, pg = new()
    base(s, "Parte 2 · 12 · Ops", "Un runbook real: sincronizar máquinas", page=pg)
    panel_bullets(s, 0.7, 2.0, 5.85, 2.7, "Outbound — copia completa",
                  ["Un tarball: workspace + ~/.claude · .aws · .ssh",
                   "-h dereferencia el symlink de .aws (crítico)",
                   "Excluye venvs / node_modules / caches / .codegraph",
                   "USB: mount manual en WSL + verificar byte a byte",
                   "En destino: bootstrap (/kg) + target-setup.sh rehacen el tooling"], accent=BLUE)
    panel_bullets(s, 6.8, 2.0, 5.85, 2.7, "Inbound — solo delta",
                  ["El código ya está en GitHub → git fetch",
                   "Docs gitignored de data/ + memoria (snapshot-memory) viajan",
                   "En la principal: restore-memory → /kg-refresh reconstruye",
                   "backup + diff de STATUS.md; STOP si hubo ediciones propias"], accent=GREEN)
    panel_bullets(s, 0.7, 4.82, 11.95, 1.4, "El landing lo conduce un agente — con guardrails",
                  ["Solo no-destructivo (renombrar, no borrar) · git fetch = única op de red",
                   "El humano hace push / merge; el agente prepara y reporta con evidencia (conteos, PRs)"],
                  accent=CORAL)
    takeaway(s, "La metodología no es solo para código: memoria durable, guardrails y 'el humano hace lo externo' también en ops.")

    # ---- 33. Cierre
    s, pg = new()
    base(s, "Cierre", "De asistente a sistema", page=pg)
    text(s, 0.7, 1.95, 5.85, 0.4, [[R("PARTE 1 · la herramienta", 13, CORAL, True)]])
    cards1 = [
        ("1", "Instalar + memoria", "Una línea; CLAUDE.md de dos niveles; sesiones que viajan."),
        ("2", "Contexto & caching", "El presupuesto y el descuento: lean y estable gana en ambos."),
        ("3", "MCP + skills + plugins", "Conecta tu mundo; empaqueta y distribuye flujos."),
        ("4", "Subagents & teams + automatización", "Escala el trabajo; hooks que garantizan calidad."),
    ]
    for i, (n, t, d) in enumerate(cards1):
        y = 2.35 + i * 0.95
        rect(s, 0.7, y, 5.85, 0.85, fill=PANEL, line=STROKE, rounded=True, radius=0.09)
        rect(s, 0.9, y + 0.18, 0.44, 0.44, fill=CORAL, rounded=True, radius=0.35)
        text(s, 0.9, y + 0.19, 0.44, 0.42, [[R(n, 13, BG, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, 1.5, y + 0.1, 5.0, 0.35, [[R(t, 12.5, TEXT, True)]])
        text(s, 1.5, y + 0.43, 5.0, 0.38, [[R(d, 10, MUTED, False)]], line_spacing=1.0)
    text(s, 6.8, 1.95, 5.85, 0.4, [[R("PARTE 2 · el método", 13, BLUE, True)]])
    cards2 = [
        ("5", "El flujo de 11 etapas", "Gates deterministas; el humano posee las decisiones."),
        ("6", "Las tools del método", "CodeGraph · Serena · GSD · /kg: barato→caro, determinista→LLM."),
        ("7", "Transferible", "Starter-kit probado en GitHub Copilot: la disciplina viaja."),
        ("8", "Hasta en ops", "Machine-sync con guardrails: el mismo método fuera del código."),
    ]
    for i, (n, t, d) in enumerate(cards2):
        y = 2.35 + i * 0.95
        rect(s, 6.8, y, 5.85, 0.85, fill=PANEL, line=STROKE, rounded=True, radius=0.09)
        rect(s, 7.0, y + 0.18, 0.44, 0.44, fill=BLUE, rounded=True, radius=0.35)
        text(s, 7.0, y + 0.19, 0.44, 0.42, [[R(n, 13, BG, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, 7.6, y + 0.1, 5.0, 0.35, [[R(t, 12.5, TEXT, True)]])
        text(s, 7.6, y + 0.43, 5.0, 0.38, [[R(d, 10, MUTED, False)]], line_spacing=1.0)
    text(s, 0.7, 6.35, 12, 0.4,
         [[R("Referencias:  code.claude.com/docs  ·  github.com/tomascortereal/claude-code-setup  ·  "
             "colbymchenry.github.io/codegraph  ·  github.com/oraios/serena", 11, FAINT, False)]])

    out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Claude_Code_Presentacion.pptx")
    prs.save(out)
    print(f"OK  ->  {out}  ({len(prs.slides._sldIdLst)} slides)")
    return out


if __name__ == "__main__":
    build()
